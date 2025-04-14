import fitz
import base64
import os
from docx import Document
from pptx import Presentation
from image_Handling import extract_text_from_image64base
import io

# fileT='./files/test.pdf'
def extract_text_from_pdf(file_bytes):
    # Create a PDF document from the bytes
    pdf_document = fitz.open(stream=file_bytes, filetype="pdf")
    text_content=[]
    # Loop through each page
    for i in range(len(pdf_document)):
        page = pdf_document.load_page(i) 
        text_content.append(page.get_text("text"))
    # Join all page texts into a single string 
    pdf_document.close()
    text = "\n".join(text_content) 
    return text

def extract_text_from_txtfile(file_bytes): 
    file_like_object = io.BytesIO(file_bytes)
    with io.TextIOWrapper(file_like_object, encoding='utf-8') as filetxt:
        text = filetxt.read()
        return text
    
def extract_text_from_docx(file):
    # Open the docx file using docx
    doc = Document(file) 
    text = []
    # Loop through each paragraph
    for paragraph in doc.paragraphs: 
        text.append(paragraph.text) 
    return "\n".join(text)

def extract_images_from_pdf(file_bytes): 
    # Create a PDF document from the bytes
    pdf_document = fitz.open(stream=file_bytes, filetype="pdf") 
    images = [] 
    # Loop through each page
    for page_num in range(len(pdf_document)): 
        page = pdf_document.load_page(page_num) 
        images_list = page.get_images(full=True) 
        for img in enumerate(images_list, start=1): 
            xref = img[0] 
            base_image = pdf_document.extract_image(xref) 
            image_bytes = base_image["image"] 
            # Convert the image bytes to base64
            base64_image = base64.b64encode(image_bytes).decode('utf-8')
            images.append(base64_image)
    pdf_document.close()
    if not images:
        return None
    return images

# Function to get the file type
def get_file_type(file_name):
    # Get the file extension
    _, ext = os.path.splitext(file_name)

    
    # Map extensions to file types
    ALLOWEDd_EXTENSIONS = {
        '.pdf': 'pdf',
        '.docx': 'docx',
        '.txt': 'txt',
        '.pptx': 'pptx',
    }
    
    # Return the corresponding file type or None if not found
    return ALLOWEDd_EXTENSIONS.get(ext, None)

def extract_text_from_pptx(file_bytes):
    # Convert the bytes to a file-like object
    file_like_object = io.BytesIO(file_bytes) 
    # Extract the text from the pptx file
    slide_text = []
    # Open the pptx file using python-pptx
    presentation = Presentation(file_like_object)
    # Loop through each slide
    for slide_num, slide in enumerate(presentation.slides, start=1):
        # Loop through each shape on the slide
        for shape in slide.shapes:
            # Check if the shape has text
            if hasattr(shape, "text"):
                # Add the text to the list
                slide_text.append(shape.text)
                # Join the text from all slides into a single string
    text = "\n".join(slide_text)
    return text           


# the main function of handling files whiche should i import in the main file :
def file_handler(file_bytes,file_name):
    """
    Handle different file types and extract text.

    Parameters:
    file_bytes (bytes): The content of the file in bytes.
    file_name (str): The name of the file including its extension.

    Returns:
    str: The extracted text from the file or an error message if the file type is not supported.
    """
    file_type = get_file_type(file_name)
    if file_type == 'pdf':
        ex_text = extract_text_from_pdf(file_bytes)
        final_text = ex_text
        ex_images = extract_images_from_pdf(file_bytes)
        if ex_images is not None:
            for img in ex_images:
                img_text = extract_text_from_image64base(img, type="png")
                final_text += "\n" + img_text
        return str(final_text)
    elif file_type == 'docx':
        ex_text = extract_text_from_docx(file_bytes)
        return str(ex_text)
    elif file_type == 'txt':
        ex_text=extract_text_from_txtfile(file_bytes)
        return str(ex_text)
    elif file_type == 'pptx':
        ex_text = extract_text_from_pptx(file_bytes)
        return str(ex_text)
    else:
        return 'file type not supported'
